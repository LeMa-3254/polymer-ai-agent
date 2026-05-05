#!/usr/bin/env python3
"""
Daily AI Tools for Polymer Materials Research Agent
----------------------------------------------------
Searches the web for the latest AI tools & research in polymer materials,
then generates a Word document and PowerPoint slide deck.

Usage:
    python3 agent.py
    ANTHROPIC_API_KEY=sk-... python3 agent.py
"""

import json
import os
import sys
from datetime import datetime
from pathlib import Path

import anthropic
from duckduckgo_search import DDGS

from docx import Document
from docx.shared import Pt, RGBColor, Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import docx.opc.constants

from pptx import Presentation
from pptx.util import Inches, Pt as PPTXPt, Emu
from pptx.dml.color import RGBColor as PPTXColor
from pptx.enum.text import PP_ALIGN


# ─────────────────────────────────────────────────────────────────────────────
# Tools
# ─────────────────────────────────────────────────────────────────────────────

def web_search(query: str, max_results: int = 6) -> list:
    try:
        with DDGS() as ddgs:
            return list(ddgs.text(query, max_results=max_results))
    except Exception as e:
        return [{"error": str(e)}]


def news_search(query: str, max_results: int = 6) -> list:
    try:
        with DDGS() as ddgs:
            return list(ddgs.news(query, max_results=max_results))
    except Exception as e:
        return [{"error": str(e)}]


TOOLS = [
    {
        "name": "web_search",
        "description": (
            "Search the web for information about AI tools and machine learning "
            "applications in polymer materials science. Returns snippets with title, "
            "body text, and URL."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Search query string"},
                "max_results": {"type": "integer", "default": 6},
            },
            "required": ["query"],
        },
    },
    {
        "name": "news_search",
        "description": (
            "Search for recent news articles about AI in polymer materials. "
            "Returns recent news items with title, body, date, and URL."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "News search query"},
                "max_results": {"type": "integer", "default": 6},
            },
            "required": ["query"],
        },
    },
]


def dispatch_tool(name: str, inputs: dict) -> str:
    if name == "web_search":
        return json.dumps(web_search(**inputs), ensure_ascii=False)
    if name == "news_search":
        return json.dumps(news_search(**inputs), ensure_ascii=False)
    return json.dumps({"error": f"Unknown tool: {name}"})


# ─────────────────────────────────────────────────────────────────────────────
# Agent loop
# ─────────────────────────────────────────────────────────────────────────────

SYSTEM_PROMPT = """You are an expert research analyst specializing in AI applications for polymer materials science.

Your task:
1. Search the web broadly for the LATEST AI tools, models, and methods used in polymer materials (prefer content from the last 6 months).
2. Cover these categories: property prediction, generative/inverse design, characterization automation, processing optimization, sustainability/recycling, LLMs in materials science.
3. Perform at least 5 different searches to gather comprehensive coverage.
4. Synthesize a structured report.

After all searches, output ONLY a single JSON object (no markdown fences, no extra text):

{
  "date": "YYYY-MM-DD",
  "headline": "one punchy sentence summarizing today's landscape",
  "executive_summary": "3-4 paragraph narrative overview",
  "findings": [
    {
      "category": "Property Prediction | Generative Design | Characterization | Processing | Sustainability | LLMs",
      "title": "concise finding title",
      "description": "what it is, what it does, and why it matters for polymer engineers (2-4 sentences)",
      "source": "URL or publication name",
      "impact": "high | medium | low"
    }
  ],
  "key_trends": ["trend 1", "trend 2", "trend 3", "trend 4", "trend 5"],
  "notable_tools": [
    {"name": "tool name", "purpose": "1-sentence purpose", "url": "URL or empty string"}
  ],
  "outlook": "Forward-looking paragraph (3-5 sentences) on where AI+polymers is heading"
}

Aim for 6-10 findings and 4-6 notable tools. Be specific — name actual tools, models, and papers."""


def run_agent(client: anthropic.Anthropic) -> dict:
    print("Starting polymer AI research agent...\n")

    messages = [
        {
            "role": "user",
            "content": (
                f"Today is {datetime.now().strftime('%Y-%m-%d')}. "
                "Research the latest AI tools and developments for polymer materials science "
                "and return a JSON report."
            ),
        }
    ]

    iteration = 0
    while True:
        iteration += 1
        response = client.messages.create(
            model="claude-opus-4-7",
            max_tokens=8000,
            system=SYSTEM_PROMPT,
            tools=TOOLS,
            messages=messages,
        )

        messages.append({"role": "assistant", "content": response.content})

        if response.stop_reason == "end_turn":
            for block in response.content:
                if hasattr(block, "text"):
                    text = block.text.strip()
                    # Strip any accidental code fences
                    if "```" in text:
                        parts = text.split("```")
                        for part in parts:
                            part = part.strip()
                            if part.startswith("json"):
                                part = part[4:].strip()
                            try:
                                return json.loads(part)
                            except json.JSONDecodeError:
                                pass
                    return json.loads(text)
            raise ValueError("No text block found in final response")

        if response.stop_reason == "tool_use":
            tool_results = []
            for block in response.content:
                if block.type == "tool_use":
                    q = block.input.get("query", "")
                    print(f"  [{iteration}] {block.name}: {q}")
                    result = dispatch_tool(block.name, block.input)
                    tool_results.append({
                        "type": "tool_result",
                        "tool_use_id": block.id,
                        "content": result,
                    })
            messages.append({"role": "user", "content": tool_results})


# ─────────────────────────────────────────────────────────────────────────────
# Word document
# ─────────────────────────────────────────────────────────────────────────────

NAVY_RGB   = RGBColor(0x0a, 0x16, 0x28)
TEAL_RGB   = RGBColor(0x0e, 0x74, 0x90)
MUTED_RGB  = RGBColor(0x64, 0x74, 0x8b)
LIGHT_RGB  = RGBColor(0x94, 0xa3, 0xb8)


def _add_hyperlink(paragraph, text, url, size_pt=9):
    """Add a clickable hyperlink; falls back to plain text if URL is missing/invalid."""
    if not url or not url.startswith("http"):
        r = paragraph.add_run(text)
        r.font.size = Pt(size_pt)
        r.font.color.rgb = LIGHT_RGB
        return
    part = paragraph.part
    r_id = part.relate_to(url, docx.opc.constants.RELATIONSHIP_TYPE.HYPERLINK, is_external=True)
    hyperlink = OxmlElement('w:hyperlink')
    hyperlink.set(qn('r:id'), r_id)
    new_run = OxmlElement('w:r')
    rPr = OxmlElement('w:rPr')
    color_el = OxmlElement('w:color')
    color_el.set(qn('w:val'), '0563C1')
    rPr.append(color_el)
    u = OxmlElement('w:u')
    u.set(qn('w:val'), 'single')
    rPr.append(u)
    sz = OxmlElement('w:sz')
    sz.set(qn('w:val'), str(size_pt * 2))
    rPr.append(sz)
    new_run.append(rPr)
    t = OxmlElement('w:t')
    t.text = text
    new_run.append(t)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)


def _heading(doc, text, size=16, color=NAVY_RGB, bold=True):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold = bold
    r.font.size = Pt(size)
    r.font.color.rgb = color
    return p


def generate_word(data: dict, path: Path):
    doc = Document()

    # Page margins
    for section in doc.sections:
        section.left_margin   = DocxInches(1.0)
        section.right_margin  = DocxInches(1.0)
        section.top_margin    = DocxInches(1.0)
        section.bottom_margin = DocxInches(1.0)

    # Title
    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    tr = title_p.add_run("AI Tools for Polymer Materials")
    tr.bold = True
    tr.font.size = Pt(22)
    tr.font.color.rgb = NAVY_RGB

    sub_p = doc.add_paragraph()
    sub_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sr = sub_p.add_run(f"Daily Research Digest  ·  {data['date']}")
    sr.font.size = Pt(11)
    sr.font.color.rgb = MUTED_RGB

    doc.add_paragraph()

    # Headline
    hl_p = doc.add_paragraph()
    hl_r = hl_p.add_run(f"📌  {data['headline']}")
    hl_r.bold = True
    hl_r.font.size = Pt(12)
    hl_r.font.color.rgb = TEAL_RGB

    doc.add_paragraph()

    # Executive Summary
    _heading(doc, "Executive Summary", size=15)
    doc.add_paragraph(data["executive_summary"])
    doc.add_paragraph()

    # Key Findings
    _heading(doc, "Key Findings", size=15)
    for i, f in enumerate(data.get("findings", []), 1):
        # Finding title
        ft = doc.add_paragraph()
        fr = ft.add_run(f"{i}.  {f['title']}")
        fr.bold = True
        fr.font.size = Pt(12)
        fr.font.color.rgb = TEAL_RGB

        # Meta line
        meta_p = doc.add_paragraph()
        meta_r = meta_p.add_run(
            f"Category: {f['category']}   |   Impact: {f['impact'].upper()}"
        )
        meta_r.italic = True
        meta_r.font.size = Pt(9)
        meta_r.font.color.rgb = MUTED_RGB

        doc.add_paragraph(f["description"])

        if f.get("source"):
            src_p = doc.add_paragraph()
            src_r = src_p.add_run("Source: ")
            src_r.font.size = Pt(9)
            src_r.font.color.rgb = LIGHT_RGB
            _add_hyperlink(src_p, f["source"], f["source"])

        doc.add_paragraph()

    # Key Trends
    _heading(doc, "Key Trends", size=15)
    for trend in data.get("key_trends", []):
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(trend)
    doc.add_paragraph()

    # Notable Tools
    tools = data.get("notable_tools", [])
    if tools:
        _heading(doc, "Notable Tools & Resources", size=15)
        for tool in tools:
            p = doc.add_paragraph(style="List Bullet")
            br = p.add_run(f"{tool['name']}: ")
            br.bold = True
            p.add_run(tool["purpose"])
            if tool.get("url"):
                p.add_run("  ")
                _add_hyperlink(p, tool["url"], tool["url"])
        doc.add_paragraph()

    # Outlook
    _heading(doc, "Outlook", size=15)
    doc.add_paragraph(data["outlook"])
    doc.add_paragraph()

    # Footer
    footer_p = doc.add_paragraph()
    footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_r = footer_p.add_run(
        f"Generated by Polymer AI Research Agent  ·  {data['date']}"
    )
    footer_r.font.size = Pt(8)
    footer_r.font.color.rgb = LIGHT_RGB

    doc.save(path)
    print(f"  Word document: {path}")


# ─────────────────────────────────────────────────────────────────────────────
# PowerPoint
# ─────────────────────────────────────────────────────────────────────────────

C_NAVY  = PPTXColor(0x0a, 0x16, 0x28)
C_TEAL  = PPTXColor(0x0e, 0x74, 0x90)
C_CYAN  = PPTXColor(0x22, 0xd3, 0xee)
C_WHITE = PPTXColor(0xFF, 0xFF, 0xFF)
C_SLATE = PPTXColor(0x64, 0x74, 0x8b)
C_LIGHT = PPTXColor(0x94, 0xa3, 0xb8)
C_BG    = PPTXColor(0xf8, 0xfa, 0xfc)
C_TEXT  = PPTXColor(0x1e, 0x29, 0x3b)
C_MTEXT = PPTXColor(0x47, 0x55, 0x69)
C_BORDER= PPTXColor(0xe2, 0xe8, 0xf0)

IMPACT_COLORS = {
    "high":   PPTXColor(0x16, 0xa3, 0x4a),
    "medium": C_TEAL,
    "low":    C_SLATE,
}


def _rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _txt(slide, text, l, t, w, h,
         size=12, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PPTXPt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _nav_bar(slide, prs, title_text):
    _rect(slide, 0, 0, prs.slide_width / 914400, prs.slide_height / 914400, C_NAVY)
    # Narrow teal bar at top
    _rect(slide, 0, 0, prs.slide_width / 914400, 0.06, C_TEAL)
    _txt(slide, title_text, 0.35, 0.12, 9.3, 0.7,
         size=22, bold=True, color=C_WHITE)


def _light_slide(slide, prs):
    _rect(slide, 0, 0, prs.slide_width / 914400, prs.slide_height / 914400, C_BG)


def _header_bar(slide, prs, label):
    _rect(slide, 0, 0, prs.slide_width / 914400, 1.0, C_NAVY)
    _rect(slide, 0, 0, prs.slide_width / 914400, 0.06, C_TEAL)
    _txt(slide, label, 0.4, 0.18, 9.2, 0.75, size=22, bold=True, color=C_WHITE)


def slide_cover(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width / 914400
    H = prs.slide_height / 914400

    _rect(slide, 0, 0, W, H, C_NAVY)
    _rect(slide, 0, 0, W, 0.06, C_TEAL)
    _rect(slide, 0, H * 0.55, W, 0.04, C_TEAL)

    _txt(slide, "AI Tools for Polymer Materials",
         0.5, 1.4, 9.0, 1.1, size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(slide, "Daily Research Digest",
         0.5, 2.6, 9.0, 0.6, size=18, color=C_CYAN, align=PP_ALIGN.CENTER)
    _txt(slide, data["date"],
         0.5, 3.2, 9.0, 0.45, size=13, color=C_SLATE, align=PP_ALIGN.CENTER)
    _txt(slide, f'"{data["headline"]}"',
         0.9, 4.1, 8.2, 0.9, size=13, bold=True, italic=True,
         color=C_LIGHT, align=PP_ALIGN.CENTER)


def slide_summary(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width / 914400
    _light_slide(slide, prs)
    _header_bar(slide, prs, "Executive Summary")

    summary = data["executive_summary"]
    if len(summary) > 700:
        summary = summary[:700].rsplit(" ", 1)[0] + "…"

    _txt(slide, summary, 0.5, 1.15, 9.0, 4.6, size=12, color=C_TEXT)

    # Trend pills
    trends = data.get("key_trends", [])[:4]
    x = 0.4
    for trend in trends:
        label = (trend[:38] + "…") if len(trend) > 38 else trend
        W_pill = 2.2
        _rect(slide, x, 6.05, W_pill, 0.42, C_TEAL)
        _txt(slide, label, x + 0.1, 6.07, W_pill - 0.15, 0.38, size=8, color=C_WHITE)
        x += W_pill + 0.17


def slide_findings(prs, data):
    findings = data.get("findings", [])
    chunks = [findings[i:i+3] for i in range(0, len(findings), 3)]

    for idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        W = prs.slide_width / 914400
        _light_slide(slide, prs)
        label = "Key Findings" if idx == 0 else "Key Findings (continued)"
        _header_bar(slide, prs, label)

        for ci, f in enumerate(chunk):
            top = 1.1 + ci * 2.0

            # Card
            _rect(slide, 0.3, top, 9.4, 1.8, C_WHITE, C_BORDER)
            # Impact stripe
            ic = IMPACT_COLORS.get(f["impact"].lower(), C_TEAL)
            _rect(slide, 0.3, top, 0.12, 1.8, ic)

            _txt(slide, f["title"], 0.55, top + 0.1, 8.0, 0.42,
                 size=11, bold=True, color=C_NAVY)
            _txt(slide, f"{f['category']}  ·  Impact: {f['impact'].upper()}",
                 0.55, top + 0.48, 8.0, 0.3, size=8, italic=True, color=C_SLATE)

            desc = f["description"]
            if len(desc) > 250:
                desc = desc[:250].rsplit(" ", 1)[0] + "…"
            _txt(slide, desc, 0.55, top + 0.76, 8.7, 0.9, size=9, color=C_MTEXT)


def slide_tools(prs, data):
    tools = data.get("notable_tools", [])[:6]
    if not tools:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width / 914400
    _light_slide(slide, prs)
    _header_bar(slide, prs, "Notable Tools & Resources")

    cols, col_w, gap = 2, 4.55, 0.2
    for i, tool in enumerate(tools):
        col = i % cols
        row = i // cols
        l = 0.3 + col * (col_w + gap)
        t = 1.1 + row * 2.0

        _rect(slide, l, t, col_w, 1.8, C_WHITE, C_BORDER)
        _rect(slide, l, t, col_w, 0.08, C_TEAL)

        _txt(slide, tool["name"], l + 0.15, t + 0.15, col_w - 0.3, 0.42,
             size=11, bold=True, color=C_NAVY)

        purpose = tool["purpose"]
        if len(purpose) > 160:
            purpose = purpose[:160].rsplit(" ", 1)[0] + "…"
        _txt(slide, purpose, l + 0.15, t + 0.55, col_w - 0.3, 1.1,
             size=9, color=C_MTEXT)


def slide_outlook(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width / 914400
    H = prs.slide_height / 914400
    _rect(slide, 0, 0, W, H, C_NAVY)
    _rect(slide, 0, 0, W, 0.06, C_TEAL)

    _txt(slide, "Outlook & Trends", 0.5, 0.5, 9.0, 0.8,
         size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    outlook = data["outlook"]
    if len(outlook) > 500:
        outlook = outlook[:500].rsplit(" ", 1)[0] + "…"
    _txt(slide, outlook, 0.7, 1.45, 8.6, 2.2,
         size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Trend list
    _rect(slide, 0, 3.8, W, 0.03, C_TEAL)
    _txt(slide, "Key Trends", 0.5, 3.95, 9.0, 0.45,
         size=13, bold=True, color=C_CYAN)

    for i, trend in enumerate(data.get("key_trends", [])[:5]):
        t = 4.45 + i * 0.44
        _rect(slide, 0.5, t + 0.12, 0.15, 0.15, C_TEAL)
        _txt(slide, trend, 0.8, t, 8.5, 0.42, size=11, color=C_WHITE)

    _txt(slide, f"Generated by Polymer AI Research Agent  ·  {data['date']}",
         0, 7.1, W, 0.3, size=8, color=C_SLATE, align=PP_ALIGN.CENTER)


def generate_pptx(data: dict, path: Path):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_cover(prs, data)
    slide_summary(prs, data)
    slide_findings(prs, data)
    slide_tools(prs, data)
    slide_outlook(prs, data)

    prs.save(path)
    print(f"  PowerPoint:    {path}")


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def main():
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        print("ERROR: ANTHROPIC_API_KEY environment variable not set.")
        print("  Run:  export ANTHROPIC_API_KEY=sk-ant-...")
        sys.exit(1)

    client = anthropic.Anthropic(api_key=api_key)

    output_dir = Path(__file__).parent / "output"
    output_dir.mkdir(exist_ok=True)
    date_str = datetime.now().strftime("%Y-%m-%d")

    # Run the research agent
    data = run_agent(client)
    n = len(data.get("findings", []))
    print(f"\nFound {n} findings, {len(data.get('key_trends', []))} trends, "
          f"{len(data.get('notable_tools', []))} tools\n")

    # Save raw JSON (useful for debugging / re-running document generation)
    json_path = output_dir / f"polymer_ai_{date_str}.json"
    json_path.write_text(json.dumps(data, indent=2, ensure_ascii=False))
    print(f"  JSON data:     {json_path}\n")

    print("Generating documents...")
    generate_word(data, output_dir / f"Polymer_AI_Digest_{date_str}.docx")
    generate_pptx(data,  output_dir / f"Polymer_AI_Digest_{date_str}.pptx")

    print(f"\nDone! All files in: {output_dir.resolve()}")

    # Open output folder
    os.system(f"open '{output_dir}'")


if __name__ == "__main__":
    main()
